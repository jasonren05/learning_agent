from flask import Flask, request, jsonify, send_file
from flask_sqlalchemy import SQLAlchemy
from flask_cors import CORS
from flask_jwt_extended import JWTManager, jwt_required, create_access_token, get_jwt_identity
from werkzeug.security import generate_password_hash, check_password_hash
from werkzeug.utils import secure_filename
import os
from datetime import datetime, timedelta
import json
from dotenv import load_dotenv
import base64
import mimetypes

# AI相关导入
from openai import OpenAI
import PyPDF2
import docx
from pptx import Presentation
import requests
from io import BytesIO
from PIL import Image

# 加载环境变量
load_dotenv()

app = Flask(__name__)
app.config['SECRET_KEY'] = os.getenv('FLASK_SECRET_KEY', 'your-secret-key')
app.config['SQLALCHEMY_DATABASE_URI'] = os.getenv('DATABASE_URL', 'sqlite:///learning_assistant.db')
app.config['SQLALCHEMY_TRACK_MODIFICATIONS'] = False
app.config['JWT_SECRET_KEY'] = os.getenv('JWT_SECRET_KEY', 'jwt-secret-string')
app.config['JWT_ACCESS_TOKEN_EXPIRES'] = timedelta(days=7)
app.config['UPLOAD_FOLDER'] = 'uploads'

# 确保上传文件夹存在
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)

db = SQLAlchemy(app)
cors = CORS(app)
jwt = JWTManager(app)

# 初始化AI客户端
ai_client = OpenAI(
    base_url=os.getenv('ARK_BASE_URL', 'https://ark.cn-beijing.volces.com/api/v3'),
    api_key=os.getenv('ARK_API_KEY')
)

# 数据库模型
class User(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    username = db.Column(db.String(80), unique=True, nullable=False)
    email = db.Column(db.String(120), unique=True, nullable=False)
    password_hash = db.Column(db.String(120), nullable=False)
    english_level = db.Column(db.String(20), default='high_school')  # 英语水平
    created_at = db.Column(db.DateTime, default=datetime.utcnow)
    
    notes = db.relationship('Note', backref='user', lazy=True)
    progress_records = db.relationship('ProgressRecord', backref='user', lazy=True)

class Note(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    title = db.Column(db.String(200), nullable=False)
    content = db.Column(db.Text)
    file_path = db.Column(db.String(500))
    file_type = db.Column(db.String(50))
    category = db.Column(db.String(100))
    keywords = db.Column(db.Text)  # JSON格式存储关键词
    created_at = db.Column(db.DateTime, default=datetime.utcnow)
    updated_at = db.Column(db.DateTime, default=datetime.utcnow, onupdate=datetime.utcnow)

class ProgressRecord(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    note_id = db.Column(db.Integer, db.ForeignKey('note.id'), nullable=False)
    access_count = db.Column(db.Integer, default=0)
    mastery_level = db.Column(db.Integer, default=0)  # 0-100的掌握程度
    last_accessed = db.Column(db.DateTime, default=datetime.utcnow)
    
    note = db.relationship('Note', backref='progress_records')

class VocabularyRecord(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    word = db.Column(db.String(100), nullable=False)
    known = db.Column(db.Boolean, default=False)
    created_at = db.Column(db.DateTime, default=datetime.utcnow)

class EnhancedContent(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    original_content = db.Column(db.Text)
    enhanced_content = db.Column(db.Text)
    content_type = db.Column(db.String(50))  # 'note', 'problem', 'english'
    is_image = db.Column(db.Boolean, default=False)
    file_path = db.Column(db.String(500))  # 保存优化内容的文件路径
    created_at = db.Column(db.DateTime, default=datetime.utcnow)
    
    user = db.relationship('User', backref='enhanced_contents')

# 文件处理工具函数
def extract_text_from_pdf(file_path):
    text = ""
    with open(file_path, 'rb') as file:
        pdf_reader = PyPDF2.PdfReader(file)
        for page in pdf_reader.pages:
            text += page.extract_text()
    return text

def extract_text_from_docx(file_path):
    doc = docx.Document(file_path)
    text = ""
    for paragraph in doc.paragraphs:
        text += paragraph.text + "\n"
    return text

def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def encode_image_to_base64(file_path):
    """将图片编码为base64"""
    try:
        with open(file_path, 'rb') as f:
            image_data = f.read()
            base64_image = base64.b64encode(image_data).decode('utf-8')
            mime_type = mimetypes.guess_type(file_path)[0]
            return f"data:{mime_type};base64,{base64_image}"
    except Exception as e:
        print(f"[ERROR] 图片编码失败: {e}")
        return None

def process_uploaded_file(file_path, file_type):
    """处理上传的文件并提取文本内容"""
    print(f"[FILE] 处理文件: {file_path}, 类型: {file_type}")
    
    if file_type in ['jpg', 'jpeg', 'png', 'gif', 'bmp', 'webp']:
        # 图片文件，返回base64编码
        return encode_image_to_base64(file_path)
    elif file_type == 'pdf':
        return extract_text_from_pdf(file_path)
    elif file_type in ['docx', 'doc']:
        return extract_text_from_docx(file_path)
    elif file_type in ['pptx', 'ppt']:
        return extract_text_from_pptx(file_path)
    elif file_type == 'txt':
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    return ""

# AI助手功能
def call_ai_api(messages, model_id=None):
    """调用豆包AI API"""
    try:
        if not model_id:
            model_id = os.getenv('ARK_MODEL_ID', 'doubao-seed-1-6-250615')
        
        print(f"[AI API] 调用模型: {model_id}")
        print(f"[AI API] 消息数量: {len(messages)}")
        
        response = ai_client.chat.completions.create(
            model=model_id,
            messages=messages
        )
        
        result = response.choices[0].message.content
        print(f"[AI API] 响应长度: {len(result) if result else 0} 字符")
        return result
    except Exception as e:
        print(f"[ERROR] AI API调用错误: {e}")
        print(f"[ERROR] 错误类型: {type(e).__name__}")
        import traceback
        print(f"[ERROR] 详细错误信息:\n{traceback.format_exc()}")
        return None

def enhance_notes(content, is_image=False):
    """补全笔记功能"""
    print(f"[ENHANCE] 开始补全笔记，是否为图片: {is_image}")
    
    if is_image:
        messages = [
            {
                "role": "system",
                "content": "你是一个专业的学习助手。请分析图片中的内容，提取关键信息，并补全完善相关知识点，使内容更加完整和易于理解。请用Markdown格式输出。"
            },
            {
                "role": "user",
                "content": [
                    {
                        "type": "text",
                        "text": "请分析这张图片中的内容，提取关键信息并补全相关知识点："
                    },
                    {
                        "type": "image_url",
                        "image_url": {
                            "url": content
                        }
                    }
                ]
            }
        ]
    else:
        messages = [
            {
                "role": "system",
                "content": "你是一个专业的学习助手。请帮助用户补全和完善笔记内容，填补知识点间的逻辑关系，使笔记更加完整和易于理解。请用Markdown格式输出，包含适当的标题、列表和重点标记。"
            },
            {
                "role": "user", 
                "content": f"请帮我补全和完善以下笔记内容，填补缺失的知识点和逻辑关系：\n\n{content}"
            }
        ]
    return call_ai_api(messages)

def generate_problem_analysis(problems, is_image=False):
    """生成题目详细解析"""
    print(f"[ANALYSIS] 开始解析题目，是否为图片: {is_image}")
    
    if is_image:
        messages = [
            {
                "role": "system",
                "content": "你是一个专业的教学助手。请分析图片中的题目，提供详细的解析，包括解题思路、步骤和知识点说明。请用Markdown格式输出。"
            },
            {
                "role": "user",
                "content": [
                    {
                        "type": "text",
                        "text": "请分析这张图片中的题目并提供详细解析："
                    },
                    {
                        "type": "image_url",
                        "image_url": {
                            "url": problems
                        }
                    }
                ]
            }
        ]
    else:
        messages = [
            {
                "role": "system",
                "content": "你是一个专业的教学助手。请为用户提供详细的题目解析，包括解题思路、步骤和知识点说明。请用Markdown格式输出。"
            },
            {
                "role": "user",
                "content": f"请为以下题目生成详细的解析：\n\n{problems}"
            }
        ]
    return call_ai_api(messages)

def save_enhanced_content(user_id, original_content, enhanced_content, content_type, is_image=False):
    """保存优化后的内容到文件和数据库"""
    try:
        # 创建保存目录
        save_dir = os.path.join('enhanced_content', str(user_id))
        os.makedirs(save_dir, exist_ok=True)
        
        # 生成文件名
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        filename = f"{content_type}_{timestamp}.md"
        file_path = os.path.join(save_dir, filename)
        
        # 保存到文件
        with open(file_path, 'w', encoding='utf-8') as f:
            f.write(f"# {content_type.title()} - 优化内容\n\n")
            f.write(f"**创建时间**: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n\n")
            if not is_image:
                f.write(f"## 原始内容\n\n{original_content}\n\n")
            f.write(f"## 优化后内容\n\n{enhanced_content}\n")
        
        # 保存到数据库
        enhanced_record = EnhancedContent(
            user_id=user_id,
            original_content=original_content if not is_image else "图片内容",
            enhanced_content=enhanced_content,
            content_type=content_type,
            is_image=is_image,
            file_path=file_path
        )
        
        db.session.add(enhanced_record)
        db.session.commit()
        
        print(f"[SAVE] 优化内容已保存: {file_path}")
        return enhanced_record.id
        
    except Exception as e:
        print(f"[ERROR] 保存优化内容失败: {e}")
        return None

def generate_english_study_material(text, user_level, is_image=False):
    """生成英语学习材料"""
    print(f"[ENGLISH] 生成英语学习材料，用户水平: {user_level}, 是否为图片: {is_image}")
    
    level_map = {
        'primary': '小学',
        'middle': '初中', 
        'high_school': '高中',
        'cet4': '大学四级',
        'cet6': '大学六级',
        'ielts_toefl': '雅思托福'
    }
    
    if is_image:
        messages = [
            {
                "role": "system",
                "content": f"你是一个专业的英语学习助手。用户的英语水平是{level_map.get(user_level, '高中')}。请分析图片中的英语内容并生成学习材料，用Markdown格式输出。"
            },
            {
                "role": "user",
                "content": [
                    {
                        "type": "text",
                        "text": f"请分析这张图片中的英语内容，根据我的英语水平({level_map.get(user_level, '高中')})生成学习材料，包括：1.内容导读 2.重点词汇及其释义 3.语法结构分析："
                    },
                    {
                        "type": "image_url",
                        "image_url": {
                            "url": text
                        }
                    }
                ]
            }
        ]
    else:
        messages = [
            {
                "role": "system",
                "content": f"你是一个专业的英语学习助手。用户的英语水平是{level_map.get(user_level, '高中')}。请根据用户水平生成学习材料，用Markdown格式输出。"
            },
            {
                "role": "user",
                "content": f"请为以下英语文章生成学习材料，包括：1.文章导读 2.超出用户水平的词汇及其英文释义 3.重点语法结构分析：\n\n{text}"
            }
        ]
    return call_ai_api(messages)

# API路由
@app.route('/api/register', methods=['POST'])
def register():
    data = request.get_json()
    
    if User.query.filter_by(username=data['username']).first():
        return jsonify({'error': '用户名已存在'}), 400
    
    if User.query.filter_by(email=data['email']).first():
        return jsonify({'error': '邮箱已被注册'}), 400
    
    user = User(
        username=data['username'],
        email=data['email'],
        password_hash=generate_password_hash(data['password']),
        english_level=data.get('english_level', 'high_school')
    )
    
    db.session.add(user)
    db.session.commit()
    
    access_token = create_access_token(identity=user.id)
    return jsonify({'access_token': access_token, 'user_id': user.id})

@app.route('/api/login', methods=['POST'])
def login():
    data = request.get_json()
    user = User.query.filter_by(username=data['username']).first()
    
    if user and check_password_hash(user.password_hash, data['password']):
        access_token = create_access_token(identity=user.id)
        return jsonify({'access_token': access_token, 'user_id': user.id})
    
    return jsonify({'error': '用户名或密码错误'}), 401

@app.route('/api/upload', methods=['POST'])
@jwt_required()
def upload_file():
    user_id = get_jwt_identity()
    
    print(f"[UPLOAD] 用户 {user_id} 开始上传文件")
    
    if 'file' not in request.files:
        print("[ERROR] 没有文件在请求中")
        return jsonify({'error': '没有文件'}), 400
    
    file = request.files['file']
    if file.filename == '':
        print("[ERROR] 文件名为空")
        return jsonify({'error': '没有选择文件'}), 400
    
    filename = secure_filename(file.filename)
    file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    file.save(file_path)
    
    print(f"[UPLOAD] 文件已保存: {file_path}")
    
    # 确定文件类型
    file_type = filename.split('.')[-1].lower()
    
    # 检查是否为图片
    is_image = file_type in ['jpg', 'jpeg', 'png', 'gif', 'bmp', 'webp']
    
    # 提取文件内容
    content = process_uploaded_file(file_path, file_type)
    
    # 创建笔记记录
    note = Note(
        user_id=user_id,
        title=request.form.get('title', filename),
        content=content,
        file_path=file_path,
        file_type=file_type,
        category=request.form.get('category', '未分类')
    )
    
    db.session.add(note)
    db.session.commit()
    
    print(f"[UPLOAD] 笔记创建成功，ID: {note.id}, 是否为图片: {is_image}")
    
    return jsonify({
        'message': '文件上传成功', 
        'note_id': note.id,
        'is_image': is_image,
        'file_type': file_type
    })

@app.route('/api/notes', methods=['GET'])
@jwt_required()
def get_notes():
    user_id = get_jwt_identity()
    notes = Note.query.filter_by(user_id=user_id).all()
    
    return jsonify([{
        'id': note.id,
        'title': note.title,
        'category': note.category,
        'created_at': note.created_at.isoformat(),
        'file_type': note.file_type
    } for note in notes])

@app.route('/api/notes/<int:note_id>', methods=['GET'])
@jwt_required()
def get_note(note_id):
    user_id = get_jwt_identity()
    note = Note.query.filter_by(id=note_id, user_id=user_id).first()
    
    if not note:
        return jsonify({'error': '笔记不存在'}), 404
    
    # 更新访问记录
    progress = ProgressRecord.query.filter_by(user_id=user_id, note_id=note_id).first()
    if not progress:
        progress = ProgressRecord(user_id=user_id, note_id=note_id, access_count=1)
        db.session.add(progress)
    else:
        progress.access_count += 1
        progress.last_accessed = datetime.utcnow()
    
    db.session.commit()
    
    return jsonify({
        'id': note.id,
        'title': note.title,
        'content': note.content,
        'category': note.category,
        'file_type': note.file_type,
        'created_at': note.created_at.isoformat()
    })

@app.route('/api/enhance-notes', methods=['POST'])
@jwt_required()
def enhance_notes_api():
    user_id = get_jwt_identity()
    data = request.get_json()
    content = data.get('content', '')
    is_image = data.get('is_image', False)
    
    print(f"[API] 用户 {user_id} 请求笔记补全，是否为图片: {is_image}")
    
    if not content:
        return jsonify({'error': '内容不能为空'}), 400
    
    enhanced_content = enhance_notes(content, is_image)
    
    if enhanced_content:
        # 保存优化后的内容
        save_id = save_enhanced_content(user_id, content, enhanced_content, 'note', is_image)
        
        return jsonify({
            'enhanced_content': enhanced_content,
            'save_id': save_id
        })
    else:
        return jsonify({'error': 'AI服务暂时不可用'}), 500

@app.route('/api/analyze-problems', methods=['POST'])
@jwt_required()
def analyze_problems():
    user_id = get_jwt_identity()
    data = request.get_json()
    problems = data.get('problems', '')
    is_image = data.get('is_image', False)
    
    print(f"[API] 用户 {user_id} 请求题目解析，是否为图片: {is_image}")
    
    if not problems:
        return jsonify({'error': '题目内容不能为空'}), 400
    
    analysis = generate_problem_analysis(problems, is_image)
    
    if analysis:
        # 保存优化后的内容
        save_id = save_enhanced_content(user_id, problems, analysis, 'problem', is_image)
        
        return jsonify({
            'analysis': analysis,
            'save_id': save_id
        })
    else:
        return jsonify({'error': 'AI服务暂时不可用'}), 500

@app.route('/api/english-study', methods=['POST'])
@jwt_required()
def english_study():
    user_id = get_jwt_identity()
    user = db.session.get(User, user_id)
    data = request.get_json()
    text = data.get('text', '')
    is_image = data.get('is_image', False)
    
    print(f"[API] 用户 {user_id} 请求英语学习材料，是否为图片: {is_image}")
    
    if not text:
        return jsonify({'error': '文章内容不能为空'}), 400
    
    study_material = generate_english_study_material(text, user.english_level, is_image)
    
    if study_material:
        # 保存优化后的内容
        save_id = save_enhanced_content(user_id, text, study_material, 'english', is_image)
        
        return jsonify({
            'study_material': study_material,
            'save_id': save_id
        })
    else:
        return jsonify({'error': 'AI服务暂时不可用'}), 500

@app.route('/api/vocabulary', methods=['POST'])
@jwt_required()
def record_vocabulary():
    user_id = get_jwt_identity()
    data = request.get_json()
    
    word = data.get('word')
    known = data.get('known', False)
    
    # 检查是否已有记录
    existing = VocabularyRecord.query.filter_by(user_id=user_id, word=word).first()
    if existing:
        existing.known = known
    else:
        record = VocabularyRecord(user_id=user_id, word=word, known=known)
        db.session.add(record)
    
    db.session.commit()
    return jsonify({'message': '记录已保存'})

@app.route('/api/progress', methods=['GET'])
@jwt_required()
def get_progress():
    user_id = get_jwt_identity()
    
    # 获取学习统计数据
    total_notes = Note.query.filter_by(user_id=user_id).count()
    total_access = db.session.query(db.func.sum(ProgressRecord.access_count)).filter_by(user_id=user_id).scalar() or 0
    avg_mastery = db.session.query(db.func.avg(ProgressRecord.mastery_level)).filter_by(user_id=user_id).scalar() or 0
    
    # 最近访问的笔记
    recent_notes = db.session.query(Note, ProgressRecord).join(
        ProgressRecord, Note.id == ProgressRecord.note_id
    ).filter(Note.user_id == user_id).order_by(
        ProgressRecord.last_accessed.desc()
    ).limit(5).all()
    
    recent_list = [{
        'title': note.title,
        'category': note.category,
        'access_count': progress.access_count,
        'last_accessed': progress.last_accessed.isoformat()
    } for note, progress in recent_notes]
    
    return jsonify({
        'total_notes': total_notes,
        'total_access': total_access,
        'average_mastery': round(avg_mastery, 1),
        'recent_notes': recent_list
    })

if __name__ == '__main__':
    with app.app_context():
        db.create_all()
    app.run(debug=True, port=5001, use_reloader=False)
